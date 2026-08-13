#!/usr/bin/env python
"""R0 자료 실측 — 비정형 문서 23종 추출 가능성 조사 스크립트.

원본은 절대 수정하지 않는다. 읽기만 하고 data/raw_probe/text/ 아래에 추출 결과를 남긴다.
페이지/슬라이드 경계는 locator 단위이므로 마커(=== [PAGE n] ===)로 보존한다.
"""

import json
import re
import subprocess
import sys
from pathlib import Path

SRC_DIR = Path("data/sources/wiki")
OUT_DIR = Path("./data/raw_probe/text")
REPORT = Path("./data/raw_probe/probe_stats.json")

# PII 탐지용 정규식 (마스킹이 아니라 "실재 여부 실측"이 목적)
RE_PHONE = re.compile(r"(?:0\d{1,2})[-\s.)]?\d{3,4}[-\s.]?\d{4}")
RE_MOBILE = re.compile(r"01[016789][-\s.]?\d{3,4}[-\s.]?\d{4}")
RE_EMAIL = re.compile(r"[A-Za-z0-9._%+\-]+@[A-Za-z0-9.\-]+\.[A-Za-z]{2,}")
RE_RRN = re.compile(r"\b\d{6}[-\s]?[1-4]\d{6}\b")


def probe_pdf(path: Path):
    """pdftotext -layout. 페이지 구분자는 form feed(\\x0c)."""
    res = subprocess.run(
        ["pdftotext", "-layout", str(path), "-"],
        capture_output=True,
    )
    raw = res.stdout.decode("utf-8", errors="replace")
    pages = raw.split("\x0c")
    if pages and not pages[-1].strip():
        pages.pop()
    chunks = []
    nonempty = 0
    for i, p in enumerate(pages, 1):
        if p.strip():
            nonempty += 1
        chunks.append(f"=== [PAGE {i}] ===\n{p}")
    text = "\n".join(chunks)
    return text, {
        "unit": "page",
        "unit_count": len(pages),
        "nonempty_units": nonempty,
        "extractor": "pdftotext -layout",
    }


def probe_pptx(path: Path):
    from pptx import Presentation

    prs = Presentation(str(path))
    chunks = []
    nonempty = 0

    def walk(shape, out):
        if shape.shape_type == 6:  # GROUP
            for sub in shape.shapes:
                walk(sub, out)
            return
        if shape.has_text_frame:
            t = shape.text_frame.text
            if t and t.strip():
                out.append(t)
        if getattr(shape, "has_table", False):
            rows = []
            for row in shape.table.rows:
                rows.append(" | ".join(c.text.replace("\n", " ") for c in row.cells))
            if rows:
                out.append("[TABLE]\n" + "\n".join(rows))

    for i, slide in enumerate(prs.slides, 1):
        out = []
        for shape in slide.shapes:
            try:
                walk(shape, out)
            except Exception as e:  # noqa: BLE001
                out.append(f"[shape error: {e}]")
        note = ""
        if slide.has_notes_slide:
            note = slide.notes_slide.notes_text_frame.text or ""
        body = "\n".join(out)
        if body.strip():
            nonempty += 1
        block = f"=== [SLIDE {i}] ===\n{body}"
        if note.strip():
            block += f"\n--- [NOTES {i}] ---\n{note}"
        chunks.append(block)
    return "\n".join(chunks), {
        "unit": "slide",
        "unit_count": len(prs.slides),
        "nonempty_units": nonempty,
        "extractor": "python-pptx",
    }


def probe_docx(path: Path):
    import docx

    d = docx.Document(str(path))
    parts = []
    para_n = 0
    for p in d.paragraphs:
        if p.text.strip():
            para_n += 1
            style = p.style.name if p.style is not None else ""
            prefix = f"[{style}] " if style and style.startswith("Heading") else ""
            parts.append(prefix + p.text)
    tbl_n = 0
    for t in d.tables:
        tbl_n += 1
        rows = []
        for row in t.rows:
            rows.append(" | ".join(c.text.replace("\n", " ") for c in row.cells))
        parts.append(f"=== [TABLE {tbl_n}] ===\n" + "\n".join(rows))
    return "\n".join(parts), {
        "unit": "paragraph+table",
        "unit_count": para_n + tbl_n,
        "nonempty_units": para_n + tbl_n,
        "paragraphs": para_n,
        "tables": tbl_n,
        "extractor": "python-docx",
    }


def _html_to_text(soup):
    for tag in soup(["script", "style"]):
        tag.decompose()
    lines = [ln.strip() for ln in soup.get_text("\n").splitlines()]
    return "\n".join(ln for ln in lines if ln)


def probe_html(path: Path):
    """bs4. Google Sites 저장본은 본문이 escape 된 채 JS 데이터 블록 안에 들어 있어
    직파싱하면 거의 안 나온다. 그럴 때는 unescape 후 재파싱한 결과를 쓴다."""
    import html as html_mod

    from bs4 import BeautifulSoup

    raw = path.read_bytes().decode("utf-8", errors="replace")
    soup = BeautifulSoup(raw, "html.parser")
    text = _html_to_text(BeautifulSoup(raw, "html.parser"))
    mode = "direct"

    unesc = raw.replace("\\u003c", "<").replace("\\u003e", ">").replace('\\"', '"')
    unesc = html_mod.unescape(unesc)
    text2 = _html_to_text(BeautifulSoup(unesc, "html.parser"))
    if len(text2) > len(text) * 1.5:
        text, mode = text2, "unescape-then-parse (embedded escaped markup)"

    sections = soup.find_all(["section", "article"])
    headings = soup.find_all(["h1", "h2", "h3"])
    return text, {
        "unit": "section/heading",
        "unit_count": len(sections) or len(headings),
        "nonempty_units": len(text.splitlines()),
        "sections": len(sections),
        "headings": len(headings),
        "img_tags": len(soup.find_all("img")),
        "extractor": f"bs4 ({mode})",
    }


def probe_md(path: Path):
    text = path.read_text(encoding="utf-8", errors="replace")
    heads = [ln for ln in text.splitlines() if ln.startswith("#")]
    return text, {
        "unit": "heading",
        "unit_count": len(heads),
        "nonempty_units": len([ln for ln in text.splitlines() if ln.strip()]),
        "extractor": "plain read",
    }


HANDLERS = {
    ".pdf": probe_pdf,
    ".pptx": probe_pptx,
    ".docx": probe_docx,
    ".html": probe_html,
    ".md": probe_md,
}


def main():
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    targets = sorted(
        p for p in SRC_DIR.iterdir()
        if p.is_file() and p.suffix.lower() in HANDLERS
    )
    only = sys.argv[1:] if len(sys.argv) > 1 else None
    results = []
    for path in targets:
        if only and not any(o in path.name for o in only):
            continue
        try:
            text, meta = HANDLERS[path.suffix.lower()](path)
        except Exception as e:  # noqa: BLE001
            results.append({
                "file": path.name, "ext": path.suffix.lower(),
                "error": f"{type(e).__name__}: {e}", "chars": 0,
            })
            print(f"FAIL {path.name}: {e}")
            continue
        out = OUT_DIR / (path.name + ".txt")
        out.write_text(text, encoding="utf-8")
        stripped = re.sub(r"=== \[[A-Z]+ \d+\] ===", "", text)
        chars = len(stripped.strip())
        nows = len(re.sub(r"\s", "", stripped))
        rec = {
            "file": path.name,
            "ext": path.suffix.lower(),
            "bytes": path.stat().st_size,
            "chars": chars,
            "chars_no_ws": nows,
            "pii": {
                "phone": len(set(RE_PHONE.findall(text))),
                "mobile": len(set(RE_MOBILE.findall(text))),
                "email": len(set(RE_EMAIL.findall(text))),
                "rrn": len(set(RE_RRN.findall(text))),
            },
            **meta,
        }
        if meta.get("unit_count"):
            rec["chars_per_unit"] = round(nows / meta["unit_count"], 1)
        results.append(rec)
        print(json.dumps(rec, ensure_ascii=False))
    REPORT.write_text(json.dumps(results, ensure_ascii=False, indent=2), encoding="utf-8")
    print(f"\n{len(results)} files -> {REPORT}")


if __name__ == "__main__":
    main()
