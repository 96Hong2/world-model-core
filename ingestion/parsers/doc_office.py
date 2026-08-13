"""DOCX · PPTX 파서.

DOCX 는 문단과 표를 본문 순서 그대로 읽는다. python-docx 의 `paragraphs` / `tables` 를 따로 돌면
표가 어느 절에 속했는지 잃어버리기 때문에 body 자식 요소를 직접 순회한다.

locator 는 제목 경로 `§Ⅱ. 현황 및 문제점>1. 카카오톡 소통 리스크`.
표는 그 뒤에 `>TABLE 3` 을 붙여 원문의 표 번호를 그대로 남긴다.

PPTX 는 슬라이드가 좌표다(`slide 17`). 도형 트리를 깊이 우선으로 훑어 그룹 안 도형과 표까지
읽고, 발표자 노트는 본문 뒤에 `[노트]` 로 붙인다. 슬라이드 순서는 파일 순서를 그대로 쓴다.
"""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

import docx
from docx.document import Document as DocxDocument
from docx.table import Table
from docx.text.paragraph import Paragraph
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from ingestion.parsers.doc_common import (
    Unit,
    build_evidences,
    build_source_record,
    normalize_text,
)

MAX_HEADING_LEVEL = 3


@dataclass(frozen=True)
class Block:
    kind: str  # paragraph | table
    index: int  # 종류별 1-based 번호 (표 번호는 원문 표기와 같다)
    text: str
    heading_path: tuple[str, ...]
    style: str = ""


def _heading_level(style_name: str) -> int | None:
    if not style_name or not style_name.startswith("Heading"):
        return None
    try:
        return int(style_name.split()[-1])
    except ValueError:
        return None


def _table_text(table: Table) -> str:
    rows = []
    for row in table.rows:
        rows.append(" | ".join(cell.text.replace("\n", " ").strip() for cell in row.cells))
    return "\n".join(rows)


def read_docx_blocks(path: Path) -> list[Block]:
    document: DocxDocument = docx.Document(str(path))
    body = document.element.body
    blocks: list[Block] = []
    heading_path: list[str] = []
    para_no = 0
    table_no = 0

    para_iter = iter(document.paragraphs)
    table_iter = iter(document.tables)

    for child in body.iterchildren():
        tag = child.tag.split("}")[-1]
        if tag == "p":
            paragraph: Paragraph = next(para_iter)
            text = paragraph.text.strip()
            if not text:
                continue
            style = paragraph.style.name if paragraph.style is not None else ""
            level = _heading_level(style)
            if level is not None and level <= MAX_HEADING_LEVEL:
                del heading_path[level - 1 :]
                heading_path.append(text)
            para_no += 1
            blocks.append(
                Block(
                    kind="paragraph",
                    index=para_no,
                    text=text,
                    heading_path=tuple(heading_path),
                    style=style,
                )
            )
        elif tag == "tbl":
            table: Table = next(table_iter)
            table_no += 1
            blocks.append(
                Block(
                    kind="table",
                    index=table_no,
                    text=_table_text(table),
                    heading_path=tuple(heading_path),
                )
            )
    return blocks


def _locator(heading_path: tuple[str, ...], suffix: str | None = None) -> str:
    parts = list(heading_path) or ["본문"]
    if suffix:
        parts.append(suffix)
    return "§" + ">".join(parts)


def _units(blocks: list[Block]) -> list[Unit]:
    units: list[Unit] = []
    buffer: list[str] = []
    buffer_path: tuple[str, ...] | None = None
    used: dict[str, int] = {}

    def flush() -> None:
        nonlocal buffer, buffer_path
        if buffer and buffer_path is not None:
            units.append(_make(buffer_path, None, "\n".join(buffer)))
        buffer, buffer_path = [], None

    def _make(path: tuple[str, ...], suffix: str | None, text: str) -> Unit:
        locator = _locator(path, suffix)
        used[locator] = used.get(locator, 0) + 1
        if used[locator] > 1:
            locator = f"{locator} ({used[locator]})"
        return Unit(locator=locator, unit="section", text=normalize_text(text))

    for block in blocks:
        if block.kind == "table":
            flush()
            units.append(_make(block.heading_path, f"TABLE {block.index}", block.text))
            continue
        if buffer_path is not None and block.heading_path != buffer_path:
            flush()
        if buffer_path is None:
            buffer_path = block.heading_path
        buffer.append(block.text)
    flush()
    return [u for u in units if u.text]


# --------------------------------------------------------------------------
# PPTX
# --------------------------------------------------------------------------
@dataclass(frozen=True)
class Slide:
    number: int
    text: str
    shape_count: int


def _shape_text(shape) -> list[str]:
    """도형 하나에서 텍스트를 뽑는다. 그룹은 안쪽 도형까지 내려간다."""
    out: list[str] = []
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            out.extend(_shape_text(child))
        return out
    if getattr(shape, "has_table", False) and shape.has_table:
        for row in shape.table.rows:
            cells = [cell.text.replace("\n", " ").strip() for cell in row.cells]
            line = " | ".join(c for c in cells if c)
            if line:
                out.append(line)
        return out
    if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
        text = shape.text_frame.text.strip()
        if text:
            out.append(text)
    return out


def read_pptx_slides(path: Path) -> list[Slide]:
    presentation = Presentation(str(path))
    slides: list[Slide] = []
    for number, slide in enumerate(presentation.slides, start=1):
        lines: list[str] = []
        shapes = 0
        for shape in slide.shapes:
            shapes += 1
            lines.extend(_shape_text(shape))
        if slide.has_notes_slide:
            note = slide.notes_slide.notes_text_frame.text.strip()
            if note:
                lines.append(f"[노트] {note}")
        slides.append(
            Slide(number=number, text="\n".join(lines), shape_count=shapes)
        )
    return slides


def _slide_units(slides: list[Slide]) -> list[Unit]:
    units: list[Unit] = []
    for slide in slides:
        text = normalize_text(slide.text)
        if not text:
            continue
        units.append(
            Unit(
                locator=f"slide {slide.number}",
                unit="slide",
                text=text,
                structured={
                    "record_kind": "doc_section",
                    "섹션": f"slide {slide.number}",
                    "row": slide.number,
                },
            )
        )
    return units


# --------------------------------------------------------------------------
def parse_source(cfg: dict) -> tuple[dict, list[dict]]:
    path = Path(cfg["path"])
    if cfg["format"] == "pptx":
        slides = read_pptx_slides(path)
        evidences = build_evidences(cfg["id"], _slide_units(slides))
        nonempty = sum(1 for s in slides if s.text.strip())
        return (
            build_source_record(
                cfg,
                parser="doc_office.pptx",
                extractor="python-pptx",
                evidence_count=len(evidences),
                notes=f"슬라이드 {len(slides)}장 · 텍스트가 있는 슬라이드 {nonempty}장",
            ),
            evidences,
        )

    blocks = read_docx_blocks(path)
    evidences = build_evidences(cfg["id"], _units(blocks))
    source = build_source_record(
        cfg,
        parser="doc_office.docx",
        extractor="python-docx",
        evidence_count=len(evidences),
        notes=(
            f"문단 {sum(1 for b in blocks if b.kind == 'paragraph')}개 · "
            f"표 {sum(1 for b in blocks if b.kind == 'table')}개"
        ),
    )
    return source, evidences
